import re
from pathlib import Path

pdf_path = Path(r"C:\Users\User\Desktop\travelplanner\updatedreport_4thsem.pdf")
pptx_path = Path(r"C:\Users\User\Desktop\travelplanner\inovative project (2).pptx")

print("PDF exists:", pdf_path.exists(), pdf_path)
print("PPTX exists:", pptx_path.exists(), pptx_path)

# ---- Extract PDF text ----
import pdfplumber

pdf_text_pages = []
with pdfplumber.open(str(pdf_path)) as pdf:
    n = len(pdf.pages)
    print("PDF pages:", n)
    for i in range(min(n, 25)):
        page = pdf.pages[i]
        text = page.extract_text() or ""
        text = re.sub(r"\s+", " ", text).strip()
        pdf_text_pages.append((i + 1, text))

print("\n--- PDF PREVIEW (first 25 pages; truncated) ---")
for pno, text in pdf_text_pages:
    if not text:
        continue
    print(f"\n[Page {pno}] {text[:800]}")

# ---- Extract PPTX text ----
from pptx import Presentation

prs = Presentation(str(pptx_path))
print("\nPPTX slides:", len(prs.slides))

slides = []
for idx, slide in enumerate(prs.slides, start=1):
    chunks = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            t = re.sub(r"\s+", " ", shape.text).strip()
            if t:
                chunks.append(t)

    # Deduplicate within slide while preserving order
    seen = set()
    uniq = []
    for t in chunks:
        if t not in seen:
            seen.add(t)
            uniq.append(t)

    slides.append((idx, uniq))

print("\n--- PPTX OUTLINE (all slides) ---")
for idx, texts in slides:
    title = texts[0] if texts else "(no text)"
    print(f"\n[Slide {idx}] {title[:160]}")
    for t in texts[1:8]:
        print(" -", t[:240])
    if len(texts) > 8:
        print(f" - ... ({len(texts)-8} more text blocks)")
